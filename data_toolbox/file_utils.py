import os
import pptx
from docx import Document
import PyPDF2


def get_all_suffix_files(directory, suffix='.csv'):
    """
    Recursively traverse a specified directory and get paths of all files with a given suffix.

    Args:
        directory (str): The path of the directory to traverse.
        suffix (str): The suffix of the files to search for. Default is '.csv'.

    Returns:
        List[str]: A list of paths of all files with the specified suffix.
    """
    target_files = []

    for root, _, files in os.walk(directory):
        for file in files:
            if file.endswith(suffix):
                target_files.append(os.path.join(root, file))

    return target_files


def extract_text_from_pptx_as_list(file_path):
    """
    Extracts all text from a PowerPoint (.pptx) file and returns it as a list of strings.

    Args:
        file_path (str): The path to the PowerPoint file.

    Returns:
        Optional[List[str]]: A list of strings, each containing text from a shape, or None if an error occurs.
    """
    try:
        # Open the PowerPoint file
        presentation = pptx.Presentation(file_path)
    except Exception as e:
        print(f"Error opening PowerPoint file: {e}")
        return None

    ppt_texts = []

    try:
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    ppt_texts.append(shape.text)
    except Exception as e:
        print(f"Error extracting text: {e}")
        return None

    return ppt_texts


def read_word_document(file_path):
    """
    Reads the text content from a Word (.docx) document and returns it as a list of strings.

    Args:
        file_path (str): The path to the Word document.

    Returns:
        Optional[List[str]]: A list of strings, each containing text from a paragraph, or None if an error occurs.
    """
    try:
        # Open the Word document
        doc = Document(file_path)
    except Exception as e:
        print(f"Error opening Word document: {e}")
        return None

    text_content = []

    try:
        # Read text content
        for paragraph in doc.paragraphs:
            text_content.append(paragraph.text)
    except Exception as e:
        print(f"Error reading text content: {e}")
        return None

    return text_content


def read_pdf_document(file_path):
    """
    Reads the text content from a PDF document and returns it as a list of strings.

    Args:
        file_path (str): The path to the PDF document.

    Returns:
        Optional[List[str]]: A list of strings, each containing text from a page, or None if an error occurs.
    """
    try:
        with open(file_path, 'rb') as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            num_pages = len(pdf_reader.pages)
            text_list = [pdf_reader.pages[page_num].extract_text() for page_num in range(num_pages)]  
    except Exception as e:
        print(f"Error reading PDF document: {e}")
        return None

    return text_list